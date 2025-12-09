"""
Iterative PowerPoint Generator with VLM Evaluation
Creates 3 slide decks with different themes/layouts, evaluates them with VLM,
and iterates until a score threshold is met, then sends to human for review.
"""

import json
import os
import sys
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
from pathlib import Path

# Import existing modules
from utils.slide_generator import SlideGenerator
from utils.presentation_builder import create_presentation_from_slides_data
from utils.vlm_analyzer import VLMAnalyzer
from utils.evaluator import SlideEvaluator


class IterativePowerPointGenerator:
    """Main orchestrator for iterative PowerPoint generation with VLM evaluation"""
    
    # Theme configurations for different presentation styles
    THEMES = {
        "executive": {
            "name": "Executive Overview",
            "style": "Business-focused, high-level impact, concise bullet points",
            "color_scheme": "Professional blue",
            "layout": "Clean, minimal, emphasis on key metrics",
            "tone": "Strategic and results-oriented"
        },
        "technical": {
            "name": "Technical Deep Dive",
            "style": "Detailed methodology, architecture, implementation",
            "color_scheme": "Technical gray-blue",
            "layout": "Structured, detailed, includes diagrams",
            "tone": "Precise and analytical"
        },
        "results": {
            "name": "Results & Impact",
            "style": "Performance metrics, practical applications, visualizations",
            "color_scheme": "Data-driven green-blue",
            "layout": "Metrics-focused, charts and graphs emphasis",
            "tone": "Evidence-based and outcome-focused"
        }
    }
    
    def __init__(self, 
                 json_input_path: str,
                 output_dir: str = "iterative_output",
                 vlm_backend: str = "auto",
                 vlm_api_key: Optional[str] = None,
                 evaluator_api_key: Optional[str] = None,
                 score_threshold: float = 75.0,
                 max_iterations: int = 5):
        """
        Initialize the iterative PowerPoint generator
        
        Args:
            json_input_path: Path to JSON file with content (retrieval output)
            output_dir: Directory to save generated presentations
            vlm_backend: VLM backend to use ("auto", "local", "ollama", "gemini", "text")
            vlm_api_key: API key for VLM (if using Gemini)
            evaluator_api_key: API key for evaluator (if using OpenAI)
            score_threshold: Minimum score required to pass (0-100)
            max_iterations: Maximum number of iterations before stopping
        """
        self.json_input_path = json_input_path
        self.output_dir = output_dir
        self.score_threshold = score_threshold
        self.max_iterations = max_iterations
        
        # Initialize components
        self.slide_generator = SlideGenerator(api_key=evaluator_api_key)
        self.vlm_analyzer = VLMAnalyzer(api_key=vlm_api_key, backend=vlm_backend)
        self.evaluator = SlideEvaluator(api_key=evaluator_api_key)
        
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)
        
        # Load input JSON
        self.input_data = self._load_input_json()
        
        print(f"✓ Initialized Iterative PowerPoint Generator")
        print(f"  Input: {json_input_path}")
        print(f"  Output: {output_dir}")
        print(f"  Score Threshold: {score_threshold}")
        print(f"  Max Iterations: {max_iterations}")
    
    def _load_input_json(self) -> Dict[str, Any]:
        """Load and validate input JSON file"""
        if not os.path.exists(self.json_input_path):
            raise FileNotFoundError(f"Input JSON file not found: {self.json_input_path}")
        
        with open(self.json_input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Validate structure
        if 'relevant_chunks' not in data and 'metadata' not in data:
            raise ValueError("Invalid JSON structure: missing 'relevant_chunks' or 'metadata'")
        
        return data
    
    def generate_three_presentations(self, iteration: int = 1) -> List[Dict[str, Any]]:
        """
        Generate 3 presentations with different themes/layouts
        
        Args:
            iteration: Current iteration number
            
        Returns:
            List of 3 presentation dictionaries with slides data
        """
        print(f"\n{'='*60}")
        print(f"Generating 3 presentations (Iteration {iteration})")
        print(f"{'='*60}")
        
        presentations = []
        
        for theme_key, theme_config in self.THEMES.items():
            print(f"\nCreating {theme_config['name']} presentation...")
            
            try:
                # Generate slides with theme-specific instructions
                slides_data = self.slide_generator.generate_slides(
                    retrieval_json_path=self.json_input_path,
                    num_slides=8,  # Adjust as needed
                    model="gpt-4o",
                    theme=theme_key  # Pass theme to generator
                )
                
                # Add theme metadata
                slides_data['theme'] = theme_key
                slides_data['theme_config'] = theme_config
                slides_data['iteration'] = iteration
                slides_data['generated_at'] = datetime.now().isoformat()
                
                presentations.append(slides_data)
                print(f"  ✓ Generated {len(slides_data.get('slides', []))} slides")
                
            except Exception as e:
                print(f"  ✗ Error generating {theme_config['name']}: {e}")
                # Continue with other themes even if one fails
                continue
        
        if len(presentations) < 3:
            print(f"\n⚠ Warning: Only {len(presentations)} presentations generated (expected 3)")
        
        return presentations
    
    def create_pptx_files(self, presentations: List[Dict[str, Any]], 
                          iteration: int = 1) -> List[str]:
        """
        Create PowerPoint files from slide data
        
        Args:
            presentations: List of presentation dictionaries
            iteration: Current iteration number
            
        Returns:
            List of paths to created PPTX files
        """
        pptx_paths = []
        
        for i, slides_data in enumerate(presentations):
            theme = slides_data.get('theme', f'theme_{i+1}')
            theme_name = self.THEMES.get(theme, {}).get('name', theme)
            
            # Create filename
            filename = f"iteration_{iteration}_{theme}_{theme_name.replace(' ', '_')}.pptx"
            output_path = os.path.join(self.output_dir, filename)
            
            try:
                create_presentation_from_slides_data(slides_data, output_path)
                pptx_paths.append(output_path)
                print(f"  ✓ Created: {filename}")
            except Exception as e:
                print(f"  ✗ Error creating PPTX for {theme_name}: {e}")
        
        return pptx_paths
    
    def evaluate_with_vlm(self, pptx_path: str) -> Dict[str, Any]:
        """
        Evaluate presentation with VLM and return score + suggestions
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            Dictionary with score, suggestions, and analysis
        """
        print(f"\nEvaluating with VLM: {os.path.basename(pptx_path)}")
        
        try:
            metadata = self.input_data.get('metadata', {})
            description = metadata.get('description', '')
            audience_type = metadata.get('audience_type', 'general')
            
            # Use VLM's enhanced evaluation method
            vlm_evaluation = self.vlm_analyzer.evaluate_with_scores(
                pptx_path,
                description=description,
                audience_type=audience_type
            )
            
            # Also get traditional evaluator scores for comparison
            slides_data = self._extract_slides_from_pptx(pptx_path)
            evaluator_scores = self.evaluator.evaluate_slides(
                slides_data=slides_data,
                retrieval_json_path=self.json_input_path,
                description=description,
                audience_type=audience_type
            )
            
            # Combine VLM and evaluator scores (prefer VLM if available, fallback to evaluator)
            overall_score = vlm_evaluation.get('overall_score', 0)
            if overall_score == 0:
                overall_score = evaluator_scores.get('overall_score', 0)
            
            # Merge scores
            combined_scores = {}
            combined_scores.update(vlm_evaluation.get('scores', {}))
            # Add evaluator scores with different keys if needed
            for key, value in evaluator_scores.get('scores', {}).items():
                if key not in combined_scores:
                    combined_scores[f"evaluator_{key}"] = value
            
            # Combine suggestions
            suggestions = self._extract_suggestions(evaluator_scores, vlm_evaluation)
            
            result = {
                "pptx_path": pptx_path,
                "success": True,
                "vlm_evaluation": vlm_evaluation,
                "evaluator_scores": evaluator_scores,
                "overall_score": overall_score,
                "scores": combined_scores,
                "suggestions": suggestions,
                "meets_threshold": overall_score >= self.score_threshold
            }
            
            print(f"  VLM Score: {vlm_evaluation.get('overall_score', 0):.1f}/100")
            print(f"  Evaluator Score: {evaluator_scores.get('overall_score', 0):.1f}/100")
            print(f"  Overall Score: {overall_score:.1f}/100")
            print(f"  Meets threshold: {result['meets_threshold']}")
            
            return result
            
        except Exception as e:
            print(f"  ✗ Error evaluating: {e}")
            import traceback
            traceback.print_exc()
            return {
                "pptx_path": pptx_path,
                "success": False,
                "error": str(e),
                "overall_score": 0,
                "meets_threshold": False
            }
    
    def _extract_slides_from_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """Extract slide data from PPTX for evaluation"""
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        slides_data = {
            "title_slide": {},
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_title = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if not slide_title and len(text) < 100:
                        slide_title = text
                    slide_text.append(text)
            
            if i == 0:
                slides_data["title_slide"] = {
                    "title": slide_title or "Presentation",
                    "subtitle": "\n".join(slide_text[1:]) if len(slide_text) > 1 else ""
                }
            else:
                slides_data["slides"].append({
                    "slide_number": i,
                    "title": slide_title or f"Slide {i}",
                    "content": slide_text
                })
        
        return slides_data
    
    def _extract_suggestions(self, evaluation: Dict[str, Any], 
                            vlm_evaluation: Dict[str, Any]) -> Dict[str, Any]:
        """Extract structured suggestions from evaluation and VLM analysis"""
        suggestions = {
            "from_evaluator": {
                "recommendations": evaluation.get('recommendations', []),
                "weaknesses": evaluation.get('weaknesses', []),
                "feedback": evaluation.get('feedback', {})
            },
            "from_vlm": {
                "suggestions": vlm_evaluation.get('suggestions', {}),
                "strengths": vlm_evaluation.get('strengths', []),
                "weaknesses": vlm_evaluation.get('weaknesses', []),
                "priority_improvements": vlm_evaluation.get('priority_improvements', [])
            },
            "priority_actions": []
        }
        
        # Create priority actions from both sources
        evaluator_recs = evaluation.get('recommendations', [])
        evaluator_weak = evaluation.get('weaknesses', [])
        vlm_improvements = vlm_evaluation.get('priority_improvements', [])
        vlm_suggestions = vlm_evaluation.get('suggestions', {})
        
        # Combine all suggestions
        all_suggestions = []
        all_suggestions.extend(evaluator_recs)
        all_suggestions.extend(evaluator_weak)
        all_suggestions.extend(vlm_improvements)
        
        # Add VLM category suggestions
        for category, items in vlm_suggestions.items():
            if isinstance(items, list):
                all_suggestions.extend(items)
        
        # Remove duplicates and prioritize
        unique_suggestions = list(dict.fromkeys(all_suggestions))  # Preserves order
        suggestions["priority_actions"] = unique_suggestions[:10]  # Top 10
        
        return suggestions
    
    def iterate_until_threshold(self) -> Dict[str, Any]:
        """
        Main iteration loop: generate, evaluate, improve until threshold is met
        
        Returns:
            Dictionary with final results and all iterations
        """
        print(f"\n{'='*60}")
        print(f"Starting Iterative PowerPoint Generation")
        print(f"Target Score: {self.score_threshold}/100")
        print(f"Max Iterations: {self.max_iterations}")
        print(f"{'='*60}\n")
        
        all_iterations = []
        best_presentation = None
        best_score = 0
        
        for iteration in range(1, self.max_iterations + 1):
            print(f"\n{'='*60}")
            print(f"ITERATION {iteration}/{self.max_iterations}")
            print(f"{'='*60}")
            
            # Step 1: Generate 3 presentations
            presentations = self.generate_three_presentations(iteration)
            
            if not presentations:
                print("✗ No presentations generated. Stopping.")
                break
            
            # Step 2: Create PPTX files
            pptx_paths = self.create_pptx_files(presentations, iteration)
            
            if not pptx_paths:
                print("✗ No PPTX files created. Stopping.")
                break
            
            # Step 3: Evaluate each presentation
            evaluations = []
            for pptx_path in pptx_paths:
                eval_result = self.evaluate_with_vlm(pptx_path)
                evaluations.append(eval_result)
                
                # Track best presentation
                score = eval_result.get('overall_score', 0)
                if score > best_score:
                    best_score = score
                    best_presentation = {
                        "pptx_path": pptx_path,
                        "slides_data": presentations[pptx_paths.index(pptx_path)],
                        "evaluation": eval_result,
                        "iteration": iteration
                    }
            
            # Store iteration results
            iteration_result = {
                "iteration": iteration,
                "presentations": presentations,
                "pptx_paths": pptx_paths,
                "evaluations": evaluations,
                "best_score": max([e.get('overall_score', 0) for e in evaluations]),
                "all_scores": [e.get('overall_score', 0) for e in evaluations]
            }
            
            all_iterations.append(iteration_result)
            
            # Check if any presentation meets threshold
            max_score = iteration_result["best_score"]
            print(f"\n{'='*60}")
            print(f"Iteration {iteration} Results:")
            print(f"  Best Score: {max_score:.1f}/100")
            print(f"  Threshold: {self.score_threshold}/100")
            print(f"{'='*60}")
            
            if max_score >= self.score_threshold:
                print(f"\n✓ SUCCESS! Score {max_score:.1f} meets threshold {self.score_threshold}")
                print(f"  Best presentation: {os.path.basename(best_presentation['pptx_path'])}")
                break
            
            if iteration < self.max_iterations:
                print(f"\n⚠ Score {max_score:.1f} below threshold. Continuing to iteration {iteration + 1}...")
                # In future iterations, we could use suggestions to improve generation
                # For now, we'll generate new variations
        
        # Final results
        final_result = {
            "success": best_score >= self.score_threshold,
            "best_score": best_score,
            "threshold": self.score_threshold,
            "iterations_completed": len(all_iterations),
            "best_presentation": best_presentation,
            "all_iterations": all_iterations,
            "ready_for_human_review": best_score >= self.score_threshold
        }
        
        return final_result
    
    def save_results(self, results: Dict[str, Any]) -> str:
        """Save final results to JSON file"""
        output_file = os.path.join(
            self.output_dir,
            f"iteration_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        )
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        
        print(f"\n✓ Results saved to: {output_file}")
        return output_file
    
    def notify_human_review(self, results: Dict[str, Any]) -> Dict[str, Any]:
        """
        Prepare notification for human review
        
        Args:
            results: Final iteration results
            
        Returns:
            Dictionary with review information
        """
        if not results.get('ready_for_human_review'):
            return {
                "status": "not_ready",
                "message": f"Score {results.get('best_score', 0):.1f} below threshold {self.score_threshold}"
            }
        
        best = results.get('best_presentation', {})
        review_info = {
            "status": "ready_for_review",
            "presentation_path": best.get('pptx_path'),
            "score": results.get('best_score'),
            "iteration": best.get('iteration'),
            "evaluation": best.get('evaluation', {}),
            "suggestions": best.get('evaluation', {}).get('suggestions', {}),
            "review_instructions": {
                "step_1": f"Review the presentation: {best.get('pptx_path')}",
                "step_2": "Check content accuracy and completeness",
                "step_3": "Verify visual design and layout",
                "step_4": "Provide feedback if changes needed",
                "step_5": "Approve or request revisions"
            }
        }
        
        # Save review info
        review_file = os.path.join(
            self.output_dir,
            "human_review_info.json"
        )
        with open(review_file, 'w', encoding='utf-8') as f:
            json.dump(review_info, f, indent=2, ensure_ascii=False)
        
        print(f"\n{'='*60}")
        print(f"READY FOR HUMAN REVIEW")
        print(f"{'='*60}")
        print(f"Presentation: {best.get('pptx_path')}")
        print(f"Score: {results.get('best_score'):.1f}/100")
        print(f"Iteration: {best.get('iteration')}")
        print(f"\nReview info saved to: {review_file}")
        print(f"{'='*60}\n")
        
        return review_info


def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Iterative PowerPoint Generator with VLM Evaluation",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic usage
  python iterative_powerpoint_generator.py retrieval_output/data.json
  
  # With custom threshold and iterations
  python iterative_powerpoint_generator.py retrieval_output/data.json --threshold 80 --max-iterations 3
  
  # With VLM backend
  python iterative_powerpoint_generator.py retrieval_output/data.json --vlm-backend gemini --vlm-api-key YOUR_KEY
        """
    )
    
    parser.add_argument('json_input', help='Path to input JSON file (retrieval output)')
    parser.add_argument('--output-dir', default='iterative_output', help='Output directory (default: iterative_output)')
    parser.add_argument('--threshold', type=float, default=75.0, help='Score threshold (0-100, default: 75.0)')
    parser.add_argument('--max-iterations', type=int, default=5, help='Maximum iterations (default: 5)')
    parser.add_argument('--vlm-backend', default='auto', choices=['auto', 'local', 'ollama', 'gemini', 'text'],
                       help='VLM backend to use (default: auto)')
    parser.add_argument('--vlm-api-key', help='VLM API key (for Gemini backend)')
    parser.add_argument('--evaluator-api-key', help='Evaluator API key (for OpenAI)')
    
    args = parser.parse_args()
    
    # Initialize generator
    generator = IterativePowerPointGenerator(
        json_input_path=args.json_input,
        output_dir=args.output_dir,
        vlm_backend=args.vlm_backend,
        vlm_api_key=args.vlm_api_key,
        evaluator_api_key=args.evaluator_api_key or os.getenv('OPENAI_API_KEY'),
        score_threshold=args.threshold,
        max_iterations=args.max_iterations
    )
    
    # Run iterations
    results = generator.iterate_until_threshold()
    
    # Save results
    generator.save_results(results)
    
    # Notify for human review
    generator.notify_human_review(results)
    
    # Print summary
    print(f"\n{'='*60}")
    print(f"SUMMARY")
    print(f"{'='*60}")
    print(f"Success: {results.get('success')}")
    print(f"Best Score: {results.get('best_score'):.1f}/100")
    print(f"Threshold: {results.get('threshold')}/100")
    print(f"Iterations: {results.get('iterations_completed')}")
    print(f"Ready for Review: {results.get('ready_for_human_review')}")
    if results.get('best_presentation'):
        print(f"Best Presentation: {results['best_presentation'].get('pptx_path')}")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    main()

