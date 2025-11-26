"""
Script to generate PowerPoint presentation from retrieval output JSON file.

Usage:
    python create_presentation.py [json_file] [output_file]
    
    Examples:
        # Use default JSON file
        python create_presentation.py
        
        # Specify JSON file
        python create_presentation.py retrieval_output/0293f39f-566a-486d-a369-fe416c4f7b64_retrieval_output.json
        
        # Specify both input and output
        python create_presentation.py input.json output.pptx

Requirements:
    pip install python-pptx
"""

import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def clean_text(text, max_length=500):
    """Clean and truncate text for slides."""
    # Remove extra whitespace and newlines
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # Remove page numbers and headers if present
    text = re.sub(r'^\d+\s*', '', text)
    
    # Extract title if present (first line that looks like a title)
    lines = text.split('\n')
    if len(lines) > 1 and len(lines[0]) < 100:
        # Might be a title
        pass
    
    # Truncate if too long
    if len(text) > max_length:
        # Try to truncate at sentence boundary
        sentences = text[:max_length].rsplit('.', 1)
        if len(sentences) > 1:
            text = sentences[0] + '.'
        else:
            text = text[:max_length] + '...'
    
    return text


def extract_title_from_text(text):
    """Extract a potential title from text."""
    lines = text.split('\n')
    for line in lines[:5]:  # Check first 5 lines
        line = line.strip()
        if line and len(line) < 150 and len(line) > 10:
            # Remove common prefixes
            line = re.sub(r'^(Abstract|I\.|II\.|III\.|IV\.|V\.|VI\.|VII\.|VIII\.|IX\.|X\.)\s*', '', line)
            if line:
                return line[:100]  # Limit title length
    return None


def create_presentation_from_json(json_path, output_path=None):
    """
    Create PowerPoint presentation from retrieval output JSON.
    
    Args:
        json_path: Path to the retrieval output JSON file
        output_path: Path for output PPTX file (default: same name as JSON with .pptx extension)
    """
    # Load JSON data
    print(f"Loading JSON file: {json_path}")
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # Extract metadata
    metadata = data.get('metadata', {})
    description = metadata.get('description', '')
    audience_type = metadata.get('audience_type', 'general')
    relevant_chunks = data.get('relevant_chunks', [])
    
    print(f"Description: {description}")
    print(f"Audience: {audience_type}")
    print(f"Relevant chunks: {len(relevant_chunks)}")
    
    # Determine number of slides (1-3 based on description)
    num_slides = 3  # Default to 3 slides
    if '1 slide' in description.lower() or 'one slide' in description.lower():
        num_slides = 1
    elif '2 slide' in description.lower() or 'two slide' in description.lower():
        num_slides = 2
    
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Organize chunks into slides
    chunks_per_slide = max(1, len(relevant_chunks) // num_slides)
    
    # Extract title from first chunk
    title_slide_text = ""
    if relevant_chunks:
        first_chunk = relevant_chunks[0]['text']
        title = extract_title_from_text(first_chunk)
        if not title:
            # Try to extract from abstract or first sentence
            lines = first_chunk.split('\n')
            for line in lines:
                if 'STRIP' in line or 'Defence' in line or 'Trojan' in line:
                    title = line.strip()[:100]
                    break
        if not title:
            title = "STRIP: A Defence Against Trojan Attacks"
        title_slide_text = title
    else:
        title = "Presentation"
    
    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title_slide_text[:80]  # Limit title length
    subtitle_shape.text = f"Audience: {audience_type.title()}\nBased on retrieved content"
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Create content slides
    for slide_idx in range(num_slides):
        start_idx = slide_idx * chunks_per_slide
        end_idx = min((slide_idx + 1) * chunks_per_slide, len(relevant_chunks))
        
        if start_idx >= len(relevant_chunks):
            break
        
        slide_chunks = relevant_chunks[start_idx:end_idx]
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"Slide {slide_idx + 1}"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add content
        content_y = Inches(1.3)
        content_height = Inches(5.5)
        content_width = Inches(9)
        
        content_box = slide.shapes.add_textbox(Inches(0.5), content_y, content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Combine chunks for this slide
        slide_text = ""
        for i, chunk in enumerate(slide_chunks):
            chunk_text = clean_text(chunk['text'], max_length=800)
            
            # Add bullet point
            if slide_text:
                slide_text += "\n\n"
            
            # Add chunk content
            slide_text += f"• {chunk_text}"
            
            # Limit total content per slide
            if len(slide_text) > 2000:
                slide_text = slide_text[:2000] + "..."
                break
        
        content_frame.text = slide_text
        
        # Format content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.2
        
        # Make first line of each bullet bold if possible
        for paragraph in content_frame.paragraphs:
            if paragraph.text.startswith('•'):
                # Try to make first sentence bold
                run = paragraph.runs[0] if paragraph.runs else None
                if run:
                    # Split at first period
                    text = paragraph.text
                    if '.' in text:
                        parts = text.split('.', 1)
                        paragraph.clear()
                        p = paragraph
                        run1 = p.add_run()
                        run1.text = parts[0] + '.'
                        run1.font.bold = True
                        run1.font.size = Pt(14)
                        if len(parts) > 1:
                            run2 = p.add_run()
                            run2.text = parts[1]
                            run2.font.size = Pt(14)
    
    # Determine output path
    if output_path is None:
        base_name = os.path.splitext(os.path.basename(json_path))[0]
        output_path = os.path.join(os.path.dirname(json_path), f"{base_name}_presentation.pptx")
    
    # Save presentation
    print(f"\nSaving presentation to: {output_path}")
    prs.save(output_path)
    print(f"✓ Presentation created successfully with {len(prs.slides)} slides!")
    
    return output_path


def main():
    """Main function to run the script."""
    import sys
    
    # Default JSON file
    json_file = "retrieval_output/0293f39f-566a-486d-a369-fe416c4f7b64_retrieval_output.json"
    
    if len(sys.argv) > 1:
        json_file = sys.argv[1]
    
    if not os.path.exists(json_file):
        print(f"Error: JSON file not found: {json_file}")
        sys.exit(1)
    
    output_file = None
    if len(sys.argv) > 2:
        output_file = sys.argv[2]
    
    try:
        output_path = create_presentation_from_json(json_file, output_file)
        print(f"\n✓ Success! Presentation saved to: {output_path}")
    except Exception as e:
        print(f"\n✗ Error creating presentation: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

