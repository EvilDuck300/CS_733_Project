"""
Generate 3 Complete Slide Decks from JSON Input
Each deck covers the same information but with different themes and layouts:
1. Executive Overview - Business-focused, high-level
2. Technical Deep Dive - Detailed methodology
3. Results & Impact - Performance metrics and applications
"""

import json
import os
import sys
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Import slide generator for AI-powered content generation
from utils.slide_generator import SlideGenerator
from utils.presentation_builder import create_presentation_from_slides_data


def extract_title_from_content(data: Dict[str, Any]) -> str:
    """Extract title from the first relevant chunk"""
    relevant_chunks = data.get('relevant_chunks', [])
    if relevant_chunks:
        first_text = relevant_chunks[0].get('text', '')
        # Try to find title (usually first line or before "Abstract")
        lines = first_text.split('\n')
        for line in lines[:5]:
            if line.strip() and len(line.strip()) < 200:
                # Likely a title if it's short and in first few lines
                if ':' in line or len(line.split()) < 15:
                    return line.strip()
        # Fallback: first non-empty line
        for line in lines:
            if line.strip() and len(line.strip()) < 200:
                return line.strip()[:100]
    return "Research Presentation"


def generate_three_presentations(json_path: str, output_dir: str = "presentations"):
    """
    Generate 3 complete slide decks from JSON input with different themes
    
    Args:
        json_path: Path to retrieval output JSON file
        output_dir: Directory to save presentations
    """
    print(f"\n{'='*60}")
    print(f"Generating 3 Themed Presentations")
    print(f"{'='*60}\n")
    
    # Load JSON data
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found: {json_path}")
        return
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    metadata = data.get('metadata', {})
    description = metadata.get('description', '')
    audience_type = metadata.get('audience_type', 'academic')
    images = data.get('images', [])
    
    # Extract title
    title = extract_title_from_content(data)
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Initialize slide generator
    slide_generator = SlideGenerator(api_key=os.getenv('OPENAI_API_KEY'))
    
    if not slide_generator.client:
        print("Warning: OpenAI API key not available. Using template-based generation.")
        # Fallback to template-based generation
        return generate_template_presentations(data, json_path, output_dir, title, images)
    
    # Generate 3 themed presentations
    themes = {
        "executive": {
            "name": "Executive Overview",
            "focus": "Business impact, high-level strategy, key metrics, ROI, practical applications",
            "style": "Concise, executive summary style, bullet points with key numbers",
            "tone": "Strategic, results-oriented, decision-focused",
            "num_slides": 8
        },
        "technical": {
            "name": "Technical Deep Dive",
            "focus": "Methodology, architecture, implementation details, technical specifications, algorithms",
            "style": "Detailed, structured, includes technical terminology and diagrams",
            "tone": "Precise, analytical, comprehensive",
            "num_slides": 10
        },
        "results": {
            "name": "Results & Impact",
            "focus": "Performance metrics, outcomes, experimental results, visualizations, practical applications",
            "style": "Data-driven, metrics-focused, includes charts/graphs emphasis",
            "tone": "Evidence-based, outcome-focused, practical",
            "num_slides": 9
        }
    }
    
    generated_presentations = []
    
    for theme_key, theme_config in themes.items():
        print(f"\n{'='*60}")
        print(f"Generating {theme_config['name']} Presentation")
        print(f"{'='*60}")
        
        try:
            # Generate slides with theme-specific instructions
            slides_data = slide_generator.generate_slides(
                retrieval_json_path=json_path,
                num_slides=theme_config['num_slides'],
                model="gpt-4o",
                theme=theme_key
            )
            
            # Add images to slides
            slides_data = add_images_to_slides(slides_data, images, theme_key)
            
            # Create PPTX file
            output_filename = f"Presentation_{theme_key}_{theme_config['name'].replace(' ', '_')}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            
            create_presentation_from_slides_data(slides_data, output_path)
            
            generated_presentations.append({
                'theme': theme_key,
                'name': theme_config['name'],
                'path': output_path,
                'slides_count': len(slides_data.get('slides', []))
            })
            
            print(f"[OK] Created: {output_filename} ({len(slides_data.get('slides', []))} slides)")
            
        except Exception as e:
            print(f"[ERROR] Error generating {theme_config['name']}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    # Summary
    print(f"\n{'='*60}")
    print(f"SUMMARY")
    print(f"{'='*60}")
    print(f"Generated {len(generated_presentations)} presentations:")
    for pres in generated_presentations:
        print(f"  • {pres['name']}: {pres['path']} ({pres['slides_count']} slides)")
    print(f"{'='*60}\n")
    
    return generated_presentations


def add_images_to_slides(slides_data: Dict[str, Any], images: List[Dict[str, Any]], 
                         theme: str) -> Dict[str, Any]:
    """Add images to appropriate slides based on theme"""
    if not images:
        return slides_data
    
    slides = slides_data.get('slides', [])
    
    # Distribute images across slides based on theme
    if theme == "executive":
        # Add key visualizations to results/metrics slides
        image_indices = [2, 4, 6]  # Add to slides 2, 4, 6
    elif theme == "technical":
        # Add diagrams to methodology/architecture slides
        image_indices = [1, 3, 5, 7]  # Add to slides 1, 3, 5, 7
    else:  # results
        # Add visualizations to results slides
        image_indices = [1, 3, 5, 7, 8]  # Add to multiple slides
    
    # Add images to slides
    image_idx = 0
    for slide_idx in image_indices:
        if slide_idx < len(slides) and image_idx < len(images):
            img = images[image_idx]
            img_path = img.get('path', '')
            
            # Convert path if needed
            if not os.path.isabs(img_path):
                # Try relative to uploads
                possible_paths = [
                    img_path,
                    os.path.join('uploads', img_path),
                    os.path.join('uploads', 'extracted_images', os.path.basename(img_path))
                ]
                for path in possible_paths:
                    if os.path.exists(path):
                        img_path = path
                        break
            
            if os.path.exists(img_path):
                if 'images' not in slides[slide_idx]:
                    slides[slide_idx]['images'] = []
                slides[slide_idx]['images'].append({
                    'path': img_path,
                    'caption': img.get('caption', f"Figure {image_idx + 1}")
                })
                image_idx += 1
    
    slides_data['slides'] = slides
    return slides_data


def extract_key_points_from_chunks(chunks: List[Dict[str, Any]], theme: str) -> List[Dict[str, Any]]:
    """Extract and organize key points from chunks based on theme"""
    key_points = []
    
    # Combine all text from relevant chunks
    all_text = "\n\n".join([chunk.get('text', '') for chunk in chunks[:8]])
    
    # Extract different types of information based on theme
    if theme == "executive":
        # Executive: Problem, Solution, Results, Impact, Conclusion
        sections = [
            ("Problem Statement", ["problem", "challenge", "issue", "vulnerability", "limitation", "security risk", "concern"]),
            ("Our Solution", ["solution", "approach", "method", "algorithm", "promsec", "introduce", "propose"]),
            ("Key Results", ["result", "performance", "accuracy", "improvement", "reduction", "achieve", "demonstrate", "97%", "order-of-magnitude"]),
            ("Business Impact", ["impact", "benefit", "cost", "efficiency", "practical", "real-world", "deployment", "integration"]),
            ("Conclusion", ["conclusion", "summary", "future", "contribution", "essential step", "trustworthiness"])
        ]
    elif theme == "technical":
        # Technical: Background, Methodology, Architecture, Implementation, Details
        sections = [
            ("Background", ["background", "introduction", "llm", "code generation", "graph neural", "ggan"]),
            ("Methodology", ["methodology", "approach", "algorithm", "ggan", "contrastive", "prompt optimization", "dual-objective"]),
            ("Architecture", ["architecture", "model", "network", "graph", "neural", "generator", "discriminator", "gcn"]),
            ("Implementation", ["implementation", "training", "optimization", "pipeline", "graph construction", "code reconstruction"]),
            ("Technical Details", ["technical", "graph", "code", "representation", "ast", "cfg", "dfg", "embedding", "similarity"]),
            ("Evaluation Framework", ["evaluation", "experiment", "dataset", "baseline", "metric", "cwe", "fuzzing"])
        ]
    else:  # results
        # Results: Overview, Setup, Performance, Comparisons, Transferability
        sections = [
            ("Overview", ["overview", "introduction", "abstract", "promsec", "objective"]),
            ("Experimental Setup", ["setup", "dataset", "baseline", "metric", "python", "java", "bandit", "spotbugs"]),
            ("Performance Results", ["result", "performance", "cwe", "accuracy", "score", "97%", "reduction", "iteration"]),
            ("Comparative Analysis", ["comparison", "baseline", "superior", "outperform", "bl1", "bl2", "promsec"]),
            ("Transferability Studies", ["transfer", "generalize", "cross", "language", "llm", "unseen", "cwe", "mask"]),
            ("Conclusion", ["conclusion", "summary", "contribution", "effectiveness", "adaptability"])
        ]
    
    # Extract content for each section
    import re
    sentence_pattern = r'(?<=[.!?])\s+'
    all_sentences = re.split(sentence_pattern, all_text)
    
    for section_title, keywords in sections:
        # Find relevant sentences
        sentences = []
        
        for sentence in all_sentences:
            sentence_lower = sentence.lower()
            # Check if sentence contains any keywords
            if any(keyword in sentence_lower for keyword in keywords):
                # Clean and truncate
                sentence = sentence.strip()
                # Remove extra whitespace
                sentence = re.sub(r'\s+', ' ', sentence)
                if len(sentence) > 30 and len(sentence) < 250:
                    sentences.append(sentence)
        
        # Take top 4-6 sentences, prioritize longer/more informative ones
        if sentences:
            # Sort by length (prefer medium-length informative sentences)
            sentences.sort(key=lambda x: abs(len(x) - 120), reverse=False)
            key_points.append({
                'title': section_title,
                'content': sentences[:6]  # Up to 6 bullet points
            })
    
    # If no key points found, create basic slides from chunks
    if not key_points:
        for i, chunk in enumerate(chunks[:6]):
            text = chunk.get('text', '')[:500]  # First 500 chars
            if text:
                # Extract first sentence as title, rest as content
                first_sentence = text.split('.')[0][:80] if '.' in text else text[:80]
                remaining = text[len(first_sentence):].strip()[:400]
                
                key_points.append({
                    'title': first_sentence if len(first_sentence) > 10 else f"Key Point {i+1}",
                    'content': [remaining] if remaining else [text[:200]]
                })
    
    return key_points


def create_fallback_slides(chunks: List[Dict[str, Any]], theme: str) -> List[Dict[str, Any]]:
    """Create slides directly from chunks if keyword extraction fails"""
    import re
    
    slides = []
    
    # Clean and split text into sentences
    for i, chunk in enumerate(chunks[:8]):  # Use up to 8 chunks
        text = chunk.get('text', '').strip()
        if not text:
            continue
        
        # Clean text
        text = re.sub(r'\s+', ' ', text)
        
        # Split into sentences
        sentences = re.split(r'(?<=[.!?])\s+', text)
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 20]
        
        if not sentences:
            continue
        
        # Use first sentence as title (or first 80 chars)
        title = sentences[0][:80] if len(sentences[0]) <= 80 else sentences[0][:77] + "..."
        
        # Use remaining sentences as content (max 5-6 per slide)
        content = sentences[1:7] if len(sentences) > 1 else sentences[:6]
        
        # Clean content items
        cleaned_content = []
        for item in content:
            item = item.strip()
            if len(item) > 20 and len(item) < 250:
                cleaned_content.append(item)
        
        if cleaned_content:
            slides.append({
                'title': title,
                'content': cleaned_content
            })
    
    # If still no slides, create at least one
    if not slides and chunks:
        first_chunk = chunks[0].get('text', '')[:1000]
        if first_chunk:
            sentences = re.split(r'(?<=[.!?])\s+', first_chunk)
            sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 20]
            
            if sentences:
                slides.append({
                    'title': "Overview" if theme == "executive" else "Introduction" if theme == "technical" else "Key Findings",
                    'content': sentences[:6]
                })
    
    return slides


def generate_template_presentations(data: Dict[str, Any], json_path: str, 
                                   output_dir: str, title: str, 
                                   images: List[Dict[str, Any]]):
    """Generate presentations from JSON content without API"""
    print("Generating presentations from JSON content (no API key required)...")
    
    relevant_chunks = data.get('relevant_chunks', [])
    metadata = data.get('metadata', {})
    description = metadata.get('description', '')
    
    presentations = []
    theme_configs = {
        "executive": {"name": "Executive Overview", "color": RGBColor(0, 51, 102)},
        "technical": {"name": "Technical Deep Dive", "color": RGBColor(64, 64, 128)},
        "results": {"name": "Results & Impact", "color": RGBColor(0, 102, 51)}
    }
    
    for theme_key, theme_info in theme_configs.items():
        print(f"\nCreating {theme_info['name']} presentation...")
        
        output_path = os.path.join(output_dir, f"Presentation_{theme_key}_{theme_info['name'].replace(' ', '_')}.pptx")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = title_slide.shapes.title
        subtitle_shape = title_slide.placeholders[1]
        
        title_shape.text = title[:80] if len(title) > 80 else title
        subtitle_shape.text = f"{theme_info['name']}\n{description[:100] if description else 'Research Presentation'}"
        
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme_info['color']
        
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
        
        # Extract key points for this theme
        key_points = extract_key_points_from_chunks(relevant_chunks, theme_key)
        
        # Debug: Check if we got key points
        print(f"  Extracted {len(key_points)} key points for {theme_key}")
        # Debug: Show first key point content
        if key_points:
            first_point = key_points[0]
            print(f"  Sample: '{first_point.get('title', 'N/A')}' with {len(first_point.get('content', []))} content items")
            if first_point.get('content'):
                print(f"    First content: {first_point['content'][0][:80]}...")
        
        if not key_points:
            print(f"  Warning: No key points extracted for {theme_key}, using fallback content extraction")
            # Fallback: create slides directly from chunks
            key_points = create_fallback_slides(relevant_chunks, theme_key)
            print(f"  Fallback created {len(key_points)} slides")
        
        # Ensure we have at least some content
        if not key_points and relevant_chunks:
            print(f"  Creating basic slides from chunks...")
            # Last resort: create simple slides from chunks
            for i, chunk in enumerate(relevant_chunks[:6]):
                text = chunk.get('text', '')[:800]
                if text:
                    import re
                    # Split into sentences
                    sentences = re.split(r'(?<=[.!?])\s+', text)
                    sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 30]
                    
                    if sentences:
                        key_points.append({
                            'title': f"Key Point {i+1}" if i == 0 else f"Section {i+1}",
                            'content': sentences[:5]
                        })
        
        # Create content slides
        slide_num = 0
        for point in key_points:
            if not point.get('content') or len(point.get('content', [])) == 0:
                print(f"  Warning: Slide '{point.get('title', 'Unknown')}' has no content, skipping")
                continue
                
            blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            title_box = blank_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = point.get('title', f"Slide {slide_num + 1}")
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = theme_info['color']
            
            # Content
            content_y = Inches(1.3)
            content_box = blank_slide.shapes.add_textbox(Inches(0.5), content_y, Inches(9), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_bottom = Inches(0.1)
            content_frame.margin_top = Inches(0.1)
            
            # Add bullet points - ensure we have content
            content_items = point.get('content', [])
            if not content_items:
                content_items = [f"Content from {point.get('title', 'section')}"]
            
            # Clear any default paragraph first
            if len(content_frame.paragraphs) > 0:
                content_frame.clear()
            
            # Clean and format content items
            import re
            
            def clean_text(text):
                """Clean text for display - fix spacing issues from PDF extraction"""
                if not text:
                    return ""
                
                text = str(text)
                # Remove non-printable characters but keep common punctuation
                text = re.sub(r'[^\x20-\x7E\u00A0-\uFFFF]', ' ', text)
                
                # More aggressive spacing fixes for PDF extraction issues
                # Add space between lowercase and uppercase (word boundaries)
                text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
                # Add space between letter and number
                text = re.sub(r'([a-zA-Z])(\d)', r'\1 \2', text)
                text = re.sub(r'(\d)([a-zA-Z])', r'\1 \2', text)
                # Add space around punctuation if missing
                text = re.sub(r'([a-zA-Z])([.,!?;:])', r'\1\2', text)  # Keep punctuation attached
                text = re.sub(r'([.,!?;:])([a-zA-Z])', r'\1 \2', text)  # Space after punctuation
                
                # Fix common concatenated words (more comprehensive list)
                common_words = ['the', 'and', 'for', 'with', 'from', 'this', 'that', 'are', 'was', 'were', 
                               'can', 'will', 'should', 'have', 'has', 'had', 'been', 'being', 'which', 'what',
                               'when', 'where', 'who', 'how', 'why', 'into', 'onto', 'upon', 'about', 'around']
                for word in common_words:
                    # Add space before common words
                    text = re.sub(r'([a-z])(' + word + r'\b)', r'\1 \2', text, flags=re.IGNORECASE)
                
                # Normalize whitespace
                text = re.sub(r'\s+', ' ', text)
                # Fix common PDF extraction issues
                text = text.replace('', '')  # Remove replacement characters
                text = text.replace('', '')  # Remove replacement characters
                text = text.strip()
                
                # Capitalize first letter of sentences
                if text:
                    text = text[0].upper() + text[1:] if len(text) > 1 else text.upper()
                
                return text
            
            # Process content items
            content_text_parts = []
            for item in content_items[:6]:  # Max 6 items per slide
                item_text = clean_text(item)
                if not item_text or len(item_text) < 15:
                    continue
                
                # Truncate intelligently at sentence boundary
                if len(item_text) > 180:
                    truncated = item_text[:180]
                    # Try to find sentence end
                    for punct in ['. ', '! ', '? ']:
                        last_punct = truncated.rfind(punct)
                        if last_punct > 80:
                            item_text = truncated[:last_punct + 1].strip()
                            break
                    else:
                        # No sentence boundary found, just truncate
                        item_text = truncated.strip() + "..."
                
                if item_text:
                    content_text_parts.append(item_text)
            
            # Set content using paragraphs
            if content_text_parts:
                # Clear existing content
                content_frame.clear()
                
                # Add each item as a separate paragraph
                for i, item_text in enumerate(content_text_parts):
                    if i == 0:
                        # Use first paragraph (created by clear())
                        p = content_frame.paragraphs[0]
                    else:
                        # Add new paragraph
                        p = content_frame.add_paragraph()
                    
                    # Set text with bullet point
                    p.text = f"• {item_text}"
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)
                    p.line_spacing = 1.2
                    p.level = 0
                    
                    # Ensure font properties are set
                    try:
                        if not p.font.name:
                            p.font.name = 'Calibri'
                        p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    except:
                        pass
            else:
                # Fallback: add at least one line with chunk content
                content_frame.clear()
                p = content_frame.paragraphs[0]
                # Try to get some text from the title or description
                fallback_text = f"Content related to: {point.get('title', 'section')}"
                p.text = f"• {fallback_text}"
                p.font.size = Pt(14)
            
            slide_num += 1
        
        # Add images if available
        if images:
            # Add one image slide per theme
            img_slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = img_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = "Key Visualizations"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = theme_info['color']
            
            # Add first available image
            for img in images[:1]:  # Add first image
                img_path = img.get('path', '')
                if not os.path.isabs(img_path):
                    possible_paths = [
                        img_path,
                        os.path.join('uploads', img_path),
                        os.path.join('uploads', 'extracted_images', os.path.basename(img_path))
                    ]
                    for path in possible_paths:
                        if os.path.exists(path):
                            img_path = path
                            break
                
                if os.path.exists(img_path):
                    try:
                        pic = img_slide.shapes.add_picture(
                            img_path, 
                            Inches(1), 
                            Inches(1.5), 
                            width=Inches(8)
                        )
                    except Exception as e:
                        print(f"  Warning: Could not add image {img_path}: {e}")
        
        # Try to save, handle file lock errors
        try:
            prs.save(output_path)
            presentations.append(output_path)
            print(f"  [OK] Created: {os.path.basename(output_path)} ({len(key_points)} content slides + title)")
        except PermissionError:
            # File might be open, try with timestamp
            import datetime
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            new_output_path = output_path.replace('.pptx', f'_{timestamp}.pptx')
            prs.save(new_output_path)
            presentations.append(new_output_path)
            print(f"  [OK] Created (with timestamp): {os.path.basename(new_output_path)} ({len(key_points)} content slides + title)")
            print(f"  Note: Original file may be open in PowerPoint")
    
    return presentations


def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Generate 3 themed slide decks from JSON input",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    
    parser.add_argument('json_input', 
                       help='Path to retrieval output JSON file',
                       default='retrieval_output/aec9ee25-a657-47bc-b85e-a9a2c1049548_retrieval_output.json',
                       nargs='?')
    parser.add_argument('--output-dir', 
                       default='presentations',
                       help='Output directory (default: presentations)')
    
    args = parser.parse_args()
    
    # Use provided JSON or default
    json_path = args.json_input
    
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found: {json_path}")
        print(f"\nUsage: python generate_three_themed_presentations.py <json_file>")
        sys.exit(1)
    
    # Generate presentations
    generate_three_presentations(json_path, args.output_dir)


if __name__ == "__main__":
    main()

