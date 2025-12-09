"""
Generate 3 Different Versions of PowerPoint Presentations
Creates 3 distinct themed presentations from the same JSON input:
1. Executive Overview - Business-focused, high-level
2. Technical Deep Dive - Detailed methodology  
3. Results & Impact - Performance metrics and applications
"""

import json
import os
import re
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def load_json_data(json_path: str) -> Dict[str, Any]:
    """Load and parse JSON file"""
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def clean_text(text: str) -> str:
    """Clean and fix text spacing issues from PDF extraction"""
    if not text:
        return ""
    
    text = str(text)
    
    # Fix encoding issues
    text = text.replace('\ufffd', '')
    text = text.replace('\u2019', "'")
    text = text.replace('\u201c', '"')
    text = text.replace('\u201d', '"')
    
    # Remove non-printable characters but keep Unicode
    text = re.sub(r'[^\x20-\x7E\u00A0-\uFFFF]', ' ', text)
    
    # Fix spacing: lowercase to uppercase
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    # Letter to number
    text = re.sub(r'([a-zA-Z])(\d)', r'\1 \2', text)
    text = re.sub(r'(\d)([a-zA-Z])', r'\1 \2', text)
    
    # Fix punctuation spacing
    text = re.sub(r'([.,!?;:])([a-zA-Z])', r'\1 \2', text)
    
    # Fix common concatenated words
    common_words = ['the', 'and', 'for', 'with', 'from', 'this', 'that', 'are', 'was', 'were',
                   'can', 'will', 'should', 'have', 'has', 'had', 'been', 'being', 'which', 'what',
                   'when', 'where', 'who', 'how', 'why', 'into', 'onto', 'upon', 'about', 'around',
                   'our', 'their', 'these', 'those', 'than', 'then', 'that', 'this']
    for word in common_words:
        text = re.sub(r'([a-z])(' + word + r'\b)', r'\1 \2', text, flags=re.IGNORECASE)
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # Capitalize first letter
    if text and len(text) > 1:
        text = text[0].upper() + text[1:]
    elif text:
        text = text.upper()
    
    return text


def extract_key_points(chunks: List[Dict[str, Any]], theme: str) -> List[Dict[str, Any]]:
    """Extract theme-specific key points from chunks"""
    key_points = []
    
    # Combine text from chunks
    all_text = "\n\n".join([chunk.get('text', '') for chunk in chunks[:10]])
    all_text = clean_text(all_text)
    
    # Split into sentences
    sentences = re.split(r'(?<=[.!?])\s+', all_text)
    clean_sentences = []
    for sent in sentences:
        sent = clean_text(sent)
        if 30 < len(sent) < 300 and len(sent.split()) > 5:
            clean_sentences.append(sent)
    
    # Define theme-specific sections
    if theme == "executive":
        sections = [
            ("Problem Statement", ["problem", "challenge", "vulnerability", "threat", "risk", "attack", "trojan"]),
            ("Our Solution", ["solution", "approach", "method", "strip", "detection", "defense", "propose"]),
            ("Key Results", ["result", "performance", "accuracy", "false", "rate", "far", "frr", "0%", "achieve"]),
            ("Business Impact", ["impact", "benefit", "cost", "efficiency", "practical", "real-world", "deployment"]),
            ("Conclusion", ["conclusion", "summary", "future", "contribution", "effective"])
        ]
    elif theme == "technical":
        sections = [
            ("Background", ["background", "introduction", "dnn", "neural network", "deep learning", "machine learning"]),
            ("Methodology", ["methodology", "approach", "algorithm", "strip", "entropy", "perturbation", "detection"]),
            ("System Design", ["system", "design", "architecture", "model", "deployment", "run-time"]),
            ("Implementation", ["implementation", "training", "dataset", "mnist", "cifar", "gtsrb"]),
            ("Technical Details", ["technical", "entropy", "distribution", "boundary", "metric", "evaluation"]),
            ("Robustness", ["robustness", "variant", "adaptive", "attack", "transparency"])
        ]
    else:  # results
        sections = [
            ("Overview", ["overview", "introduction", "abstract", "strip", "objective"]),
            ("Experimental Setup", ["setup", "dataset", "baseline", "metric", "experiment", "evaluation"]),
            ("Performance Results", ["result", "performance", "accuracy", "false", "rate", "0%", "reduction"]),
            ("Comparative Analysis", ["comparison", "baseline", "superior", "outperform", "versus"]),
            ("Transferability", ["transfer", "generalize", "cross", "language", "unseen", "variant"]),
            ("Conclusion", ["conclusion", "summary", "contribution", "effectiveness"])
        ]
    
    # Extract content for each section
    for section_title, keywords in sections:
        matching_sentences = []
        
        for sentence in clean_sentences:
            sentence_lower = sentence.lower()
            matches = sum(1 for keyword in keywords if keyword.lower() in sentence_lower)
            
            if matches > 0:
                # Score by matches and length
                score = matches * 10 + (10 - abs(len(sentence) - 120) / 10)
                matching_sentences.append((score, sentence))
        
        # Sort and take top sentences
        matching_sentences.sort(key=lambda x: x[0], reverse=True)
        selected = [sent for _, sent in matching_sentences[:6]]
        
        if selected:
            key_points.append({
                'title': section_title,
                'content': selected
            })
    
    # Fallback if no key points found
    if not key_points:
        for i, chunk in enumerate(chunks[:6]):
            text = chunk.get('text', '')[:500]
            if text:
                text = clean_text(text)
                sentences = re.split(r'(?<=[.!?])\s+', text)
                sentences = [clean_text(s) for s in sentences if 40 < len(clean_text(s)) < 250]
                
                if sentences:
                    key_points.append({
                        'title': f"Key Point {i+1}",
                        'content': sentences[:5]
                    })
    
    return key_points


def create_presentation(theme_key: str, theme_info: Dict, key_points: List[Dict], 
                      title: str, description: str, images: List[Dict], 
                      output_path: str):
    """Create a single themed presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title[:80] if len(title) > 80 else title
    subtitle_shape.text = f"{theme_info['name']}\n{description[:100] if description else 'Research Presentation'}"
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_info['color']
    
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    
    # Content slides
    for point in key_points:
        if not point.get('content'):
            continue
        
        blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = blank_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = point.get('title', 'Slide')
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = theme_info['color']
        title_frame.paragraphs[0].font.name = 'Calibri'
        
        # Content
        content_box = blank_slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.clear()
        
        # Add bullet points
        content_items = point.get('content', [])
        for i, item in enumerate(content_items[:6]):
            item_text = clean_text(item)
            if not item_text or len(item_text) < 15:
                continue
            
            # Truncate if too long
            if len(item_text) > 180:
                truncated = item_text[:180]
                for punct in ['. ', '! ', '? ']:
                    last_punct = truncated.rfind(punct)
                    if last_punct > 80:
                        item_text = truncated[:last_punct + 1].strip()
                        break
                else:
                    item_text = truncated.strip() + "..."
            
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {item_text}"
            p.font.size = Pt(14)
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(8)
            p.line_spacing = 1.2
    
    # Add image slide if available
    if images:
        img_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = img_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Key Visualizations"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = theme_info['color']
        
        # Add first available image
        for img in images[:1]:
            img_path = img.get('path', '')
            if not os.path.isabs(img_path):
                possible_paths = [
                    img_path,
                    os.path.join('uploads', img_path),
                    os.path.join('uploads', 'extracted_images', os.path.basename(img_path))
                ]
                for path in possible_paths:
                    if os.path.exists(path):
                        img_path = path
                        break
            
            if os.path.exists(img_path):
                try:
                    img_slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))
                except Exception as e:
                    print(f"  Warning: Could not add image {img_path}: {e}")
    
    # Save presentation
    prs.save(output_path)
    return output_path


def generate_three_ppt_versions(json_path: str, output_dir: str = "presentations"):
    """
    Generate 3 different themed PowerPoint presentations
    
    Args:
        json_path: Path to retrieval output JSON file
        output_dir: Directory to save presentations
    """
    print(f"\n{'='*60}")
    print(f"Generating 3 Different PPT Versions")
    print(f"{'='*60}\n")
    
    # Load JSON data
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found: {json_path}")
        return []
    
    data = load_json_data(json_path)
    
    # Extract data
    relevant_chunks = data.get('relevant_chunks', [])
    metadata = data.get('metadata', {})
    description = metadata.get('description', '')
    images = data.get('images', [])
    
    # Extract title
    if relevant_chunks:
        first_text = relevant_chunks[0].get('text', '')
        lines = first_text.split('\n')
        title = "Research Presentation"
        for line in lines[:5]:
            if line.strip() and len(line.strip()) < 200:
                if ':' in line or len(line.split()) < 15:
                    title = line.strip()
                    break
    else:
        title = "Research Presentation"
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Theme configurations
    themes = {
        "executive": {
            "name": "Executive Overview",
            "color": RGBColor(0, 51, 102)
        },
        "technical": {
            "name": "Technical Deep Dive",
            "color": RGBColor(64, 64, 128)
        },
        "results": {
            "name": "Results & Impact",
            "color": RGBColor(0, 102, 51)
        }
    }
    
    generated_files = []
    
    # Generate each themed presentation
    for theme_key, theme_info in themes.items():
        print(f"\nCreating {theme_info['name']} presentation...")
        
        # Extract key points for this theme
        key_points = extract_key_points(relevant_chunks, theme_key)
        print(f"  Extracted {len(key_points)} key points")
        
        # Create output path
        output_filename = f"Presentation_{theme_key}_{theme_info['name'].replace(' ', '_')}.pptx"
        output_path = os.path.join(output_dir, output_filename)
        
        # Create presentation
        try:
            create_presentation(theme_key, theme_info, key_points, title, description, images, output_path)
            generated_files.append(output_path)
            print(f"  [OK] Created: {output_filename} ({len(key_points)} content slides + title)")
        except Exception as e:
            print(f"  [ERROR] Failed to create {theme_info['name']}: {e}")
            import traceback
            traceback.print_exc()
    
    # Summary
    print(f"\n{'='*60}")
    print(f"SUMMARY")
    print(f"{'='*60}")
    print(f"Generated {len(generated_files)} presentations:")
    for file_path in generated_files:
        print(f"  • {os.path.basename(file_path)}")
    print(f"{'='*60}\n")
    
    return generated_files


def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Generate 3 different themed PowerPoint presentations from JSON input"
    )
    
    parser.add_argument('json_input',
                       help='Path to retrieval output JSON file',
                       default='retrieval_output/22b010da-f861-42f8-912e-18ed701b25e6_retrieval_output.json',
                       nargs='?')
    parser.add_argument('--output-dir',
                       default='presentations',
                       help='Output directory (default: presentations)')
    
    args = parser.parse_args()
    
    generate_three_ppt_versions(args.json_input, args.output_dir)


if __name__ == "__main__":
    main()

