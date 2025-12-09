"""
Script to generate three different PowerPoint presentations from retrieval output JSON file.

Each presentation covers the full content but with different focus:
1. Executive Overview - Business-focused, high-level impact
2. Technical Deep Dive - Detailed methodology and architecture
3. Results & Impact - Performance metrics and practical applications
"""

import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def clean_text(text, max_length=400):
    """Clean and truncate text for slides."""
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    text = re.sub(r'^\d+\s*', '', text)
    
    if len(text) > max_length:
        sentences = text[:max_length].rsplit('.', 1)
        if len(sentences) > 1:
            text = sentences[0] + '.'
        else:
            text = text[:max_length] + '...'
    
    return text


def add_image_to_slide(slide, img_path, x, y, width, caption=None):
    """Add an image to a slide with optional caption."""
    if not os.path.exists(img_path):
        print(f"[Warning] Image not found: {img_path}")
        return None
    
    try:
        pic = slide.shapes.add_picture(img_path, x, y, width=width)
        
        if caption:
            cap_y = y + pic.height + Inches(0.1)
            cap_box = slide.shapes.add_textbox(x, cap_y, width, Inches(0.3))
            tb = cap_box.text_frame
            tb.text = caption
            tb.paragraphs[0].font.size = Pt(10)
            tb.paragraphs[0].font.italic = True
            tb.paragraphs[0].alignment = PP_ALIGN.CENTER
            tb.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        return pic
    except Exception as e:
        print(f"[ERROR] Could not add image {img_path}: {e}")
        return None


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide."""
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)


def create_content_slide(prs, title, content_items, image_path=None, image_caption=None):
    """Create a content slide with optional image."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Determine layout based on whether we have an image
    if image_path and os.path.exists(image_path):
        content_width = Inches(5)
        content_x = Inches(0.5)
        img_x = Inches(6)
        img_y = Inches(1.3)
        img_width = Inches(3.5)
    else:
        content_width = Inches(9)
        content_x = Inches(0.5)
    
    # Add content
    content_y = Inches(1.3)
    content_height = Inches(5.5)
    
    content_box = slide.shapes.add_textbox(content_x, content_y, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Add content items
    for i, item in enumerate(content_items):
        if i > 0:
            content_frame.add_paragraph()
        
        p = content_frame.paragraphs[i] if i < len(content_frame.paragraphs) else content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    
    # Add image if provided
    if image_path:
        add_image_to_slide(slide, image_path, img_x, img_y, img_width, image_caption)
    
    return slide


def create_presentation_1_executive_overview(json_path, output_path):
    """Create Executive Overview Presentation - Business-focused, high-level."""
    print("Creating Presentation 1: Executive Overview...")
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    images = data.get('images', [])
    
    # Title slide
    create_title_slide(
        prs,
        "Automated Tow-Gap Detection in AFP Composites",
        "Executive Overview\nAhmad Ravangard, Nathan Sage, Aamr Ibrahim"
    )
    
    # The Challenge
    create_content_slide(
        prs,
        "The Challenge",
        [
            "• Automated Fiber Placement (AFP) is critical for aerospace manufacturing",
            "• Manufacturing defects reduce structural integrity and performance",
            "• Tow-gap defects are among the most critical flaws",
            "• Manual inspection is time-consuming and error-prone",
            "• Need for automated, reliable quality assurance"
        ]
    )
    
    # Our Solution
    create_content_slide(
        prs,
        "Our Solution",
        [
            "• End-to-end computer vision pipeline for automated detection",
            "• Deep learning-based segmentation with UNet architecture",
            "• Real-time inspection capability",
            "• 97% defect detection rate",
            "• Suitable for inline manufacturing integration"
        ]
    )
    
    # Key Results
    create_content_slide(
        prs,
        "Key Results",
        [
            "• 97% True Positive Rate - Detects nearly all defects",
            "• 3% False Negative Rate - Minimal missed defects",
            "• F1 Score: 0.885 - Strong segmentation accuracy",
            "• IoU: 0.794 - Precise boundary detection",
            "• Computational efficiency - Works on standard CPU hardware"
        ]
    )
    
    # Business Impact
    create_content_slide(
        prs,
        "Business Impact",
        [
            "• Reduces operator dependency and inspection time",
            "• Enables real-time quality control during production",
            "• Prevents defective components from reaching assembly",
            "• Improves manufacturing consistency and reliability",
            "• Cost-effective solution for industrial deployment"
        ]
    )
    
    # Technical Approach (Brief)
    create_content_slide(
        prs,
        "Technical Approach",
        [
            "• High-resolution image processing (512×512 pixel tiles)",
            "• Expert-annotated ground truth for training",
            "• Data augmentation for robustness",
            "• UNet with ResNet34 encoder architecture",
            "• Optimized threshold calibration (t=0.95)"
        ]
    )
    
    # Add AFP process image
    page2_img = next((img for img in images if img.get('id') == 'page2_img1'), None)
    if page2_img:
        img_path = page2_img.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Automated Fiber Placement",
                [
                    "• Robotic placement head with consolidation roller",
                    "• Layered deposition of composite tows",
                    "• Defects can arise during layup process",
                    "• Our system detects gaps in real-time"
                ],
                img_path,
                "Figure 1: AFP Process"
            )
    
    # Performance Metrics
    create_content_slide(
        prs,
        "Performance Metrics",
        [
            "• Accuracy: 98.5%",
            "• Precision: 83.5%",
            "• Recall: 94.3%",
            "• Specificity: 98.8%",
            "• Balanced Accuracy: 96.5%"
        ]
    )
    
    # Conclusion
    create_content_slide(
        prs,
        "Conclusion",
        [
            "• Proven solution for automated tow-gap detection",
            "• High accuracy with computational efficiency",
            "• Ready for industrial integration",
            "• Enables data-driven quality assurance",
            "• Reduces manufacturing variability and improves reliability"
        ]
    )
    
    prs.save(output_path)
    print(f"[OK] Presentation 1 saved: {output_path}")
    return output_path


def create_presentation_2_technical_deepdive(json_path, output_path):
    """Create Technical Deep Dive Presentation - Detailed methodology."""
    print("Creating Presentation 2: Technical Deep Dive...")
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    images = data.get('images', [])
    
    # Title slide
    create_title_slide(
        prs,
        "Automated Tow-Gap Detection",
        "Technical Deep Dive\nMethodology, Architecture, and Implementation"
    )
    
    # Problem Statement
    create_content_slide(
        prs,
        "Problem Statement",
        [
            "• AFP produces tow-gap defects during manufacturing",
            "• Defects introduce resin-rich pockets and fiber waviness",
            "• Reduces laminate mechanical performance",
            "• Requires precise detection and geometric characterization",
            "• Must work under variable lighting and surface conditions"
        ]
    )
    
    # Data Acquisition
    create_content_slide(
        prs,
        "Data Acquisition",
        [
            "• High-resolution inspection images from AFP panels",
            "• Controlled industrial lighting conditions",
            "• Images normalized and partitioned into 512×512 pixel tiles",
            "• Overlapping tiling strategy to avoid truncating defects",
            "• Captures surface variations from tow placement and resin distribution"
        ]
    )
    
    # Annotation Process
    page3_img = next((img for img in images if img.get('id') == 'page3_img1'), None)
    if page3_img:
        img_path = page3_img.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Annotation and Ground Truth",
                [
                    "• Polygon-based annotations using LabelMe tool",
                    "• Expert annotators manually traced gap contours",
                    "• JSON format converted to binary masks",
                    "• Pixel-level correspondence for supervised learning",
                    "• Both defect-containing and defect-free tiles included"
                ],
                img_path,
                "Figure 2: Dataset preparation workflow"
            )
    else:
        create_content_slide(
            prs,
            "Annotation and Ground Truth",
            [
                "• Polygon-based annotations using LabelMe tool",
                "• Expert annotators manually traced gap contours",
                "• JSON format converted to binary masks",
                "• Pixel-level correspondence for supervised learning",
                "• Both defect-containing and defect-free tiles included"
            ]
        )
    
    # Data Augmentation
    create_content_slide(
        prs,
        "Data Augmentation Strategy",
        [
            "• Rotational transformations: simulate tow orientation variations",
            "• Perspective perturbations: model oblique viewing angles",
            "• Contrast adjustments: emulate illumination fluctuations",
            "• Random application to increase dataset diversity",
            "• Maintains physical integrity of ground truth labels"
        ]
    )
    
    # Model Architecture
    create_content_slide(
        prs,
        "Model Architecture",
        [
            "• UNet segmentation network with ResNet34 encoder",
            "• Lightweight design for computational efficiency",
            "• Designed for thin, elongated tow-gap structures",
            "• Maintains sharp boundary details for metrology",
            "• Suitable for precision measurement tasks"
        ]
    )
    
    # Training Strategy
    create_content_slide(
        prs,
        "Training Strategy",
        [
            "• Class-balance normalization via defect-free masks",
            "• Extensive data augmentation for robustness",
            "• Training, validation, and test set split",
            "• Optimized for lighting, orientation, and surface variability",
            "• Threshold calibration for optimal performance"
        ]
    )
    
    # Evaluation Metrics
    create_content_slide(
        prs,
        "Evaluation Metrics",
        [
            "• Intersection-over-Union (IoU) for overlap measurement",
            "• Dice/F1 score for boundary agreement",
            "• Precision and Recall for classification",
            "• Dimensional agreement (width, length)",
            "• Confusion matrix for behavior analysis"
        ]
    )
    
    # Results Overview
    create_content_slide(
        prs,
        "Results Overview",
        [
            "• 97% true positive rate",
            "• 3% false negative rate",
            "• F1 score: 0.885",
            "• IoU: 0.794",
            "• Optimal threshold: t = 0.95"
        ]
    )
    
    # Technical Implementation
    create_content_slide(
        prs,
        "Technical Implementation",
        [
            "• Real-time inference capability",
            "• CPU-based processing (no GPU required)",
            "• Tiling inference for full spatial resolution",
            "• Sub-two-pixel boundary error at high confidence",
            "• Suitable for inline AFP production integration"
        ]
    )
    
    prs.save(output_path)
    print(f"[OK] Presentation 2 saved: {output_path}")
    return output_path


def create_presentation_3_results_impact(json_path, output_path):
    """Create Results & Impact Presentation - Performance metrics and applications."""
    print("Creating Presentation 3: Results & Impact...")
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    images = data.get('images', [])
    
    # Title slide
    create_title_slide(
        prs,
        "Automated Tow-Gap Detection",
        "Results & Impact\nPerformance Evaluation and Applications"
    )
    
    # Introduction
    create_content_slide(
        prs,
        "Overview",
        [
            "• End-to-end computer vision pipeline for tow-gap detection",
            "• UNet segmentation with ResNet34 encoder",
            "• Evaluated on AFP laminate inspection imagery",
            "• High detection reliability with computational efficiency",
            "• Suitable for near-real-time industrial inspection"
        ]
    )
    
    # Confusion Matrix Results
    page4_img = next((img for img in images if img.get('id') == 'page4_img1'), None)
    if page4_img:
        img_path = page4_img.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Classification Performance",
                [
                    "• True Positive Rate: 97% - Detects 97% of defects",
                    "• False Negative Rate: 3% - Misses only 3%",
                    "• True Negative Rate: 55%",
                    "• False Positive Rate: 45% - Conservative strategy",
                    "• Strong diagonal dominance in confusion matrix"
                ],
                img_path,
                "Figure 3: Confusion Matrix"
            )
    else:
        create_content_slide(
            prs,
            "Classification Performance",
            [
                "• True Positive Rate: 97% - Detects 97% of defects",
                "• False Negative Rate: 3% - Misses only 3%",
                "• True Negative Rate: 55%",
                "• False Positive Rate: 45% - Conservative strategy",
                "• Strong diagonal dominance in confusion matrix"
            ]
        )
    
    # Segmentation Quality
    page5_img1 = next((img for img in images if img.get('id') == 'page5_img1'), None)
    if page5_img1:
        img_path = page5_img1.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Segmentation Quality Metrics",
                [
                    "• IoU (Positive-only): 0.75-0.90 range",
                    "• F1 Score (Positive-only): Above 0.85, centered at 0.90",
                    "• Reliable localization of tow gap shapes",
                    "• Accurate boundary alignment",
                    "• Consistent across lighting and surface conditions"
                ],
                img_path,
                "Figure 4: IoU and F1 distributions"
            )
    else:
        create_content_slide(
            prs,
            "Segmentation Quality Metrics",
            [
                "• IoU (Positive-only): 0.75-0.90 range",
                "• F1 Score (Positive-only): Above 0.85, centered at 0.90",
                "• Reliable localization of tow gap shapes",
                "• Accurate boundary alignment",
                "• Consistent across lighting and surface conditions"
            ]
        )
    
    # Performance Metrics Table
    create_content_slide(
        prs,
        "Comprehensive Performance Metrics",
        [
            "• Accuracy: 0.985 (98.5%)",
            "• Precision: 0.835 (83.5%)",
            "• Recall: 0.943 (94.3%)",
            "• Specificity: 0.988 (98.8%)",
            "• F1 Score: 0.885",
            "• IoU: 0.794",
            "• Balanced Accuracy: 0.965 (96.5%)",
            "• Matthews Correlation Coefficient: 0.879"
        ]
    )
    
    # Threshold Analysis
    page5_img2 = next((img for img in images if img.get('id') == 'page5_img2'), None)
    if page5_img2:
        img_path = page5_img2.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Threshold Optimization",
                [
                    "• Optimal threshold: t = 0.95",
                    "• F1 score peaks at 0.885",
                    "• High threshold filters false positives",
                    "• Well-calibrated probability outputs",
                    "• Monotonic increase in performance with threshold"
                ],
                img_path,
                "Figure 5: Metrics vs Threshold"
            )
    else:
        create_content_slide(
            prs,
            "Threshold Optimization",
            [
                "• Optimal threshold: t = 0.95",
                "• F1 score peaks at 0.885",
                "• High threshold filters false positives",
                "• Well-calibrated probability outputs",
                "• Monotonic increase in performance with threshold"
            ]
        )
    
    # Visual Results
    page6_img = next((img for img in images if img.get('id') == 'page6_img1'), None)
    if page6_img:
        img_path = page6_img.get('path', '')
        if os.path.exists(img_path):
            create_content_slide(
                prs,
                "Visual Segmentation Results",
                [
                    "• Predicted masks (green) align with visible gaps",
                    "• Identifies elongated triangular gaps",
                    "• Detects thin linear gaps between tows",
                    "• Robust to varying orientations",
                    "• Minimal false activation in defect-free regions"
                ],
                img_path,
                "Figure 6: Final segmentation results"
            )
    else:
        create_content_slide(
            prs,
            "Visual Segmentation Results",
            [
                "• Predicted masks align with visible gaps",
                "• Identifies elongated triangular gaps",
                "• Detects thin linear gaps between tows",
                "• Robust to varying orientations",
                "• Minimal false activation in defect-free regions"
            ]
        )
    
    # Computational Efficiency
    create_content_slide(
        prs,
        "Computational Efficiency",
        [
            "• Real-time inference capability",
            "• Works on standard CPU hardware",
            "• No specialist GPU required",
            "• Suitable for inline production integration",
            "• Maintains high accuracy with efficiency"
        ]
    )
    
    # Practical Applications
    create_content_slide(
        prs,
        "Practical Applications",
        [
            "• Inline quality control during AFP production",
            "• Automated acceptance testing",
            "• Process monitoring and feedback",
            "• Geometry-aware measurements (width, length, orientation)",
            "• Targeted rework identification"
        ]
    )
    
    # Impact and Benefits
    create_content_slide(
        prs,
        "Impact and Benefits",
        [
            "• Reduces operator dependency",
            "• Accelerates defect identification",
            "• Improves manufacturing consistency",
            "• Enables data-driven quality assurance",
            "• Shifts from subjective to objective evaluation"
        ]
    )
    
    # Conclusion
    create_content_slide(
        prs,
        "Conclusion",
        [
            "• High detection reliability (97% true positive rate)",
            "• Strong segmentation accuracy (F1: 0.885, IoU: 0.794)",
            "• Computational efficiency for real-time use",
            "• Ready for industrial integration",
            "• Enables reliable, data-driven AFP manufacturing"
        ]
    )
    
    prs.save(output_path)
    print(f"[OK] Presentation 3 saved: {output_path}")
    return output_path


def main():
    """Main function to generate all three presentations."""
    json_path = "retrieval_output/cba85529-0653-4b8d-b508-be5524cd3c2e_retrieval_output.json"
    
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found: {json_path}")
        return
    
    output_dir = "retrieval_output"
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        output1 = os.path.join(output_dir, "Presentation_1_Executive_Overview.pptx")
        output2 = os.path.join(output_dir, "Presentation_2_Technical_DeepDive.pptx")
        output3 = os.path.join(output_dir, "Presentation_3_Results_Impact.pptx")
        
        create_presentation_1_executive_overview(json_path, output1)
        create_presentation_2_technical_deepdive(json_path, output2)
        create_presentation_3_results_impact(json_path, output3)
        
        print("\n" + "="*60)
        print("[SUCCESS] All three presentations created successfully!")
        print("="*60)
        print(f"1. {output1}")
        print(f"2. {output2}")
        print(f"3. {output3}")
        
    except Exception as e:
        print(f"\n[ERROR] Error creating presentations: {str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
