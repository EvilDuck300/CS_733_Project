"""
VLM (Vision Language Model) Analyzer for PowerPoint Presentations
Supports multiple VLM backends: Local models (BLIP-2, LLaVA), Ollama, Gemini API, or text-based analysis
"""

import os
import json
import base64
import re
from typing import Dict, Any, List, Optional
from pptx import Presentation
from io import BytesIO

# Check for Gemini API
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    genai = None

# Check for PIL/Pillow
try:
    from PIL import Image
    import io
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    Image = None

# Check for Hugging Face Transformers (for local models)
try:
    from transformers import BlipProcessor, BlipForConditionalGeneration, AutoProcessor, AutoModelForCausalLM
    import torch
    TRANSFORMERS_AVAILABLE = True
except ImportError:
    TRANSFORMERS_AVAILABLE = False
    BlipProcessor = None
    BlipForConditionalGeneration = None
    AutoProcessor = None
    AutoModelForCausalLM = None
    torch = None

# Check for Ollama
try:
    import requests
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False
    requests = None


class VLMAnalyzer:
    """Analyze PowerPoint presentations using Vision Language Model
    
    Supports multiple backends:
    - Local models (BLIP-2, LLaVA via Hugging Face)
    - Ollama (local models)
    - Gemini API (optional)
    - Text-based analysis (fallback)
    """
    
    def __init__(self, api_key: Optional[str] = None, 
                 backend: str = "auto",
                 model_name: Optional[str] = None):
        """
        Initialize the VLM analyzer
        
        Args:
            api_key: API key (Gemini if using 'gemini' backend)
            backend: Backend to use ("auto", "local", "ollama", "gemini", "text")
            model_name: Model name for local/Ollama (e.g., "Salesforce/blip-image-captioning-base", "llava")
        """
        self.backend = backend
        self.model_name = model_name
        self.client = None
        self.local_model = None
        self.local_processor = None
        self.api_key = None
        
        if not PIL_AVAILABLE:
            print("Warning: PIL/Pillow not available. Install with: pip install Pillow")
        
        # Auto-detect backend if "auto"
        if backend == "auto":
            backend = self._detect_backend()
        
        # Initialize selected backend
        if backend == "gemini":
            self._init_gemini(api_key)
        elif backend == "local":
            self._init_local_model(model_name)
        elif backend == "ollama":
            self._init_ollama(model_name)
        elif backend == "text":
            print("Using text-based analysis (no vision model)")
        else:
            print(f"Warning: Unknown backend '{backend}'. Using text-based analysis.")
    
    def _detect_backend(self) -> str:
        """Auto-detect available backend"""
        # Check for Gemini API key
        if GEMINI_AVAILABLE and os.getenv('GEMINI_API_KEY'):
            return "gemini"
        
        # Check for local transformers models
        if TRANSFORMERS_AVAILABLE:
            return "local"
        
        # Check for Ollama
        if OLLAMA_AVAILABLE:
            try:
                response = requests.get("http://localhost:11434/api/tags", timeout=2)
                if response.status_code == 200:
                    return "ollama"
            except:
                pass
        
        # Fallback to text-based
        return "text"
    
    def _init_gemini(self, api_key: Optional[str] = None):
        """Initialize Gemini API backend"""
        if not GEMINI_AVAILABLE:
            print("Warning: google-generativeai not installed. Install with: pip install google-generativeai")
            return
        
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        if not self.api_key:
            print("Warning: Gemini API key not provided.")
            return
        
        genai.configure(api_key=self.api_key)
        self.client = genai.GenerativeModel('gemini-pro-vision')
        self.backend = "gemini"
        print("✓ VLM Analyzer initialized with Gemini Vision API")
    
    def _init_local_model(self, model_name: Optional[str] = None):
        """Initialize local Hugging Face model"""
        if not TRANSFORMERS_AVAILABLE:
            print("Warning: transformers not installed. Install with: pip install transformers torch")
            print("Falling back to text-based analysis.")
            self.backend = "text"
            return
        
        if not PIL_AVAILABLE:
            print("Warning: PIL/Pillow required for local models. Install with: pip install Pillow")
            self.backend = "text"
            return
        
        # Default model
        if model_name is None:
            model_name = "Salesforce/blip-image-captioning-base"  # Lightweight, good for slides
        
        try:
            print(f"Loading local model: {model_name}...")
            self.local_processor = AutoProcessor.from_pretrained(model_name)
            self.local_model = AutoModelForCausalLM.from_pretrained(model_name)
            
            # Use CPU if CUDA not available
            device = "cuda" if torch and torch.cuda.is_available() else "cpu"
            self.local_model.to(device)
            self.local_model.eval()
            
            self.backend = "local"
            self.model_name = model_name
            print(f"✓ Local VLM model loaded on {device}")
        except Exception as e:
            print(f"Warning: Could not load local model: {e}")
            print("Falling back to text-based analysis.")
            self.backend = "text"
    
    def _init_ollama(self, model_name: Optional[str] = None):
        """Initialize Ollama backend"""
        if not OLLAMA_AVAILABLE:
            print("Warning: requests not installed. Install with: pip install requests")
            self.backend = "text"
            return
        
        if model_name is None:
            model_name = "llava"  # Default Ollama vision model
        
        # Check if Ollama is running
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            if response.status_code != 200:
                print("Warning: Ollama not running. Start Ollama or use another backend.")
                self.backend = "text"
                return
        except:
            print("Warning: Ollama not available. Install Ollama from https://ollama.ai/")
            self.backend = "text"
            return
        
        self.backend = "ollama"
        self.model_name = model_name
        print(f"✓ Ollama backend initialized (model: {model_name})")
    
    def pptx_to_images(self, pptx_path: str, output_dir: Optional[str] = None) -> List[str]:
        """
        Convert PowerPoint slides to images
        
        Args:
            pptx_path: Path to PowerPoint file
            output_dir: Directory to save images (default: same directory as PPTX)
            
        Returns:
            List of image file paths
        """
        if not PIL_AVAILABLE:
            raise ValueError("PIL/Pillow required for slide-to-image conversion")
        
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        # Determine output directory
        if output_dir is None:
            base_dir = os.path.dirname(pptx_path)
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            output_dir = os.path.join(base_dir, f"{base_name}_slides")
        
        os.makedirs(output_dir, exist_ok=True)
        
        # Load presentation
        prs = Presentation(pptx_path)
        image_paths = []
        
        print(f"Converting {len(prs.slides)} slides to images...")
        
        # Note: python-pptx doesn't directly support rendering to images
        # We'll need to use a workaround or external tool
        # For now, we'll create a placeholder that indicates this needs implementation
        # In production, you might use:
        # - LibreOffice headless mode
        # - comtypes (Windows COM automation)
        # - unoconv
        # - Or convert PPTX to PDF first, then PDF to images
        
        print("Note: Direct PPTX to image conversion requires additional tools.")
        print("Consider using: LibreOffice, unoconv, or convert PPTX->PDF->Images")
        
        return image_paths
    
    def analyze_slide_image(self, image_path: str, prompt: str = None) -> Dict[str, Any]:
        """
        Analyze a single slide image using the configured VLM backend
        
        Args:
            image_path: Path to slide image
            prompt: Analysis prompt (default: comprehensive slide analysis)
            
        Returns:
            Dictionary with analysis results
        """
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Image file not found: {image_path}")
        
        if not PIL_AVAILABLE:
            raise ValueError("PIL/Pillow required for image processing")
        
        # Load image
        img = Image.open(image_path)
        
        # Default prompt
        if prompt is None:
            prompt = "Describe this PowerPoint slide in detail, including content, layout, and visual design."
        
        # Route to appropriate backend
        if self.backend == "gemini":
            return self._analyze_with_gemini(img, prompt)
        elif self.backend == "local":
            return self._analyze_with_local_model(img, prompt)
        elif self.backend == "ollama":
            return self._analyze_with_ollama(img, prompt)
        else:
            return self._analyze_text_only(image_path)
    
    def _analyze_with_gemini(self, img: Image.Image, prompt: str) -> Dict[str, Any]:
        """Analyze using Gemini API"""
        if not self.client:
            raise ValueError("Gemini API not initialized")
        
        try:
            response = self.client.generate_content([prompt, img])
            analysis_text = response.text
            
            return {
                "success": True,
                "backend": "gemini",
                "analysis": self._parse_response(analysis_text)
            }
        except Exception as e:
            return {"success": False, "error": str(e), "backend": "gemini"}
    
    def _analyze_with_local_model(self, img: Image.Image, prompt: str) -> Dict[str, Any]:
        """Analyze using local Hugging Face model"""
        if not self.local_model or not self.local_processor:
            raise ValueError("Local model not initialized")
        
        try:
            # Process image and generate caption/description
            inputs = self.local_processor(images=img, text=prompt, return_tensors="pt")
            
            device = next(self.local_model.parameters()).device
            inputs = {k: v.to(device) for k, v in inputs.items()}
            
            with torch.no_grad():
                outputs = self.local_model.generate(**inputs, max_length=512)
            
            analysis_text = self.local_processor.decode(outputs[0], skip_special_tokens=True)
            
            return {
                "success": True,
                "backend": "local",
                "model": self.model_name,
                "analysis": {
                    "description": analysis_text,
                    "prompt": prompt
                }
            }
        except Exception as e:
            return {"success": False, "error": str(e), "backend": "local"}
    
    def _analyze_with_ollama(self, img: Image.Image, prompt: str) -> Dict[str, Any]:
        """Analyze using Ollama"""
        if not OLLAMA_AVAILABLE:
            raise ValueError("Ollama not available")
        
        try:
            # Convert image to base64
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            
            # Call Ollama API
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": self.model_name or "llava",
                    "prompt": prompt,
                    "images": [img_base64],
                    "stream": False
                },
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                analysis_text = result.get("response", "")
                
                return {
                    "success": True,
                    "backend": "ollama",
                    "model": self.model_name or "llava",
                    "analysis": {
                        "description": analysis_text,
                        "prompt": prompt
                    }
                }
            else:
                return {"success": False, "error": f"Ollama API error: {response.status_code}"}
        except Exception as e:
            return {"success": False, "error": str(e), "backend": "ollama"}
    
    def _analyze_text_only(self, image_path: str) -> Dict[str, Any]:
        """Fallback text-based analysis (extract text from PowerPoint)"""
        # This is a placeholder - in practice, you'd extract text from the PPTX
        return {
            "success": True,
            "backend": "text",
            "analysis": {
                "message": "Text-based analysis: Extract text from PowerPoint slides for analysis",
                "note": "For visual analysis, use a vision model backend (local, ollama, or gemini)"
            }
        }
    
    def _parse_response(self, text: str) -> Dict[str, Any]:
        """Try to parse response as JSON, fallback to text"""
        try:
            if "```json" in text:
                json_text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                json_text = text.split("```")[1].split("```")[0].strip()
            else:
                json_text = text
            
            return json.loads(json_text)
        except json.JSONDecodeError:
            return {"analysis": text, "raw_response": True}
    
    def analyze_presentation(self, pptx_path: str, 
                           analysis_type: str = "comprehensive",
                           generate_improved: bool = False) -> Dict[str, Any]:
        """
        Analyze entire PowerPoint presentation using VLM
        
        Args:
            pptx_path: Path to PowerPoint file
            analysis_type: Type of analysis ("comprehensive", "visual", "content", "quality")
            generate_improved: If True, generate improved slide content (not just extract)
            
        Returns:
            Dictionary with complete analysis results and optionally improved slides
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        # Load presentation to get metadata
        prs = Presentation(pptx_path)
        num_slides = len(prs.slides)
        
        print(f"Analyzing presentation with {num_slides} slides using {self.backend} backend...")
        
        # Extract text from slides
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_title = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Try to identify title (usually first shape or larger font)
                    if not slide_title and len(text) < 100:
                        slide_title = text
                    slide_text.append(text)
            
            slides_data.append({
                "slide_number": i + 1,
                "title": slide_title or f"Slide {i + 1}",
                "content": slide_text,
                "text": "\n".join(slide_text)
            })
        
        result = {
            "success": True,
            "backend": self.backend,
            "presentation_path": pptx_path,
            "num_slides": num_slides,
            "analysis_type": analysis_type,
            "original_slides": slides_data
        }
        
        # Generate improved slides if requested
        if generate_improved:
            print("Generating improved slide content...")
            improved_slides = self._generate_improved_slides(slides_data, analysis_type)
            result["improved_slides"] = improved_slides
            result["has_improvements"] = True
        else:
            result["message"] = f"Text extracted from {num_slides} slides. Set generate_improved=True to generate improved content."
            result["has_improvements"] = False
        
        return result
    
    def _generate_improved_slides(self, original_slides: List[Dict[str, Any]], 
                                 analysis_type: str) -> Dict[str, Any]:
        """
        Generate improved slide content based on VLM analysis
        
        Args:
            original_slides: List of original slide data
            analysis_type: Type of analysis
            
        Returns:
            Dictionary with improved slides in standard format
        """
        improved_slides = {
            "title_slide": {
                "title": original_slides[0].get("title", "Presentation") if original_slides else "Presentation",
                "subtitle": "Improved with VLM Analysis"
            },
            "slides": []
        }
        
        # Build prompt for improvement
        improvement_prompt = self._build_improvement_prompt(original_slides, analysis_type)
        
        # Process each slide
        for slide_data in original_slides:
            slide_num = slide_data.get("slide_number", 0)
            original_title = slide_data.get("title", f"Slide {slide_num}")
            original_content = slide_data.get("content", [])
            
            # Generate improved content based on backend
            if self.backend == "gemini" and self.client:
                improved = self._improve_slide_with_gemini(original_title, original_content, improvement_prompt)
            elif self.backend == "local" and self.local_model:
                improved = self._improve_slide_with_local(original_title, original_content, improvement_prompt)
            elif self.backend == "ollama":
                improved = self._improve_slide_with_ollama(original_title, original_content, improvement_prompt)
            else:
                # Fallback: use rule-based improvements
                improved = self._improve_slide_rule_based(original_title, original_content)
            
            improved_slides["slides"].append({
                "slide_number": slide_num,
                "title": improved.get("title", original_title),
                "content": improved.get("content", original_content),
                "notes": improved.get("notes", f"Improved version of original slide {slide_num}")
            })
        
        return improved_slides
    
    def _build_improvement_prompt(self, slides: List[Dict[str, Any]], analysis_type: str) -> str:
        """Build prompt for improving slides"""
        return f"""You are an expert presentation designer. Improve the following PowerPoint slides by:

1. Making content clearer and more concise
2. Improving bullet points for better readability
3. Enhancing titles to be more engaging
4. Ensuring professional language and structure
5. Maintaining all key information while improving presentation

Original slides:
{json.dumps(slides, indent=2)}

Generate improved slides with:
- Better titles
- Clearer, more concise bullet points
- Professional formatting
- All original key information preserved

Return as JSON with improved title and content array."""
    
    def _improve_slide_with_gemini(self, title: str, content: List[str], prompt: str) -> Dict[str, Any]:
        """Improve slide using Gemini API"""
        if not self.client:
            return self._improve_slide_rule_based(title, content)
        
        try:
            slide_text = f"Title: {title}\nContent:\n" + "\n".join([f"- {item}" for item in content])
            full_prompt = f"{prompt}\n\nSlide to improve:\n{slide_text}\n\nProvide improved version:"
            
            response = self.client.generate_content(full_prompt)
            result_text = response.text
            
            # Try to parse JSON response
            try:
                if "```json" in result_text:
                    json_text = result_text.split("```json")[1].split("```")[0].strip()
                elif "```" in result_text:
                    json_text = result_text.split("```")[1].split("```")[0].strip()
                else:
                    json_text = result_text
                
                improved = json.loads(json_text)
                return improved
            except json.JSONDecodeError:
                # If not JSON, extract improvements from text
                return self._extract_improvements_from_text(result_text, title, content)
        except Exception as e:
            print(f"Error improving with Gemini: {e}")
            return self._improve_slide_rule_based(title, content)
    
    def _improve_slide_with_local(self, title: str, content: List[str], prompt: str) -> Dict[str, Any]:
        """Improve slide using local model"""
        if not self.local_model:
            return self._improve_slide_rule_based(title, content)
        
        try:
            slide_text = f"Title: {title}\nContent:\n" + "\n".join([f"- {item}" for item in content])
            full_prompt = f"{prompt}\n\nSlide: {slide_text}"
            
            inputs = self.local_processor(text=full_prompt, return_tensors="pt")
            device = next(self.local_model.parameters()).device
            inputs = {k: v.to(device) for k, v in inputs.items()}
            
            with torch.no_grad():
                outputs = self.local_model.generate(**inputs, max_length=512)
            
            result_text = self.local_processor.decode(outputs[0], skip_special_tokens=True)
            return self._extract_improvements_from_text(result_text, title, content)
        except Exception as e:
            print(f"Error improving with local model: {e}")
            return self._improve_slide_rule_based(title, content)
    
    def _improve_slide_with_ollama(self, title: str, content: List[str], prompt: str) -> Dict[str, Any]:
        """Improve slide using Ollama"""
        if not OLLAMA_AVAILABLE:
            return self._improve_slide_rule_based(title, content)
        
        try:
            slide_text = f"Title: {title}\nContent:\n" + "\n".join([f"- {item}" for item in content])
            full_prompt = f"{prompt}\n\nSlide: {slide_text}"
            
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": self.model_name or "llama2",
                    "prompt": full_prompt,
                    "stream": False
                },
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                result_text = result.get("response", "")
                return self._extract_improvements_from_text(result_text, title, content)
            else:
                return self._improve_slide_rule_based(title, content)
        except Exception as e:
            print(f"Error improving with Ollama: {e}")
            return self._improve_slide_rule_based(title, content)
    
    def _extract_improvements_from_text(self, text: str, original_title: str, 
                                       original_content: List[str]) -> Dict[str, Any]:
        """Extract improved content from model response text"""
        # Try to find title and content in the response
        improved_title = original_title
        improved_content = original_content.copy()
        
        # Look for title patterns
        title_patterns = [
            r'Title[:\s]+(.+?)(?:\n|$)',
            r'Improved Title[:\s]+(.+?)(?:\n|$)',
            r'^#\s*(.+?)$'
        ]
        
        for pattern in title_patterns:
            match = re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
            if match:
                improved_title = match.group(1).strip()
                break
        
        # Look for bullet points or content
        bullet_patterns = [
            r'[-•*]\s*(.+?)(?:\n|$)',
            r'\d+\.\s*(.+?)(?:\n|$)',
        ]
        
        found_bullets = []
        for pattern in bullet_patterns:
            matches = re.findall(pattern, text, re.MULTILINE)
            if matches:
                found_bullets = [m.strip() for m in matches if len(m.strip()) > 10]
                break
        
        if found_bullets:
            improved_content = found_bullets[:len(original_content) + 2]  # Allow some extra
        
        return {
            "title": improved_title,
            "content": improved_content,
            "notes": "Generated by VLM analysis"
        }
    
    def _improve_slide_rule_based(self, title: str, content: List[str]) -> Dict[str, Any]:
        """Rule-based slide improvement (fallback)"""
        # Basic improvements without AI
        improved_title = title.strip()
        if not improved_title.endswith(('.', '!', '?')):
            improved_title = improved_title.rstrip('.')
        
        improved_content = []
        for item in content:
            item = item.strip()
            # Remove redundant words, improve formatting
            if item and len(item) > 5:
                # Ensure it starts with capital
                if item[0].islower():
                    item = item[0].upper() + item[1:]
                improved_content.append(item)
        
        return {
            "title": improved_title,
            "content": improved_content,
            "notes": "Rule-based improvements applied"
        }
    
    def analyze_with_gemini_website(self, pptx_path: str) -> Dict[str, Any]:
        """
        Provide instructions for analyzing PowerPoint using Gemini website
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            Dictionary with instructions and file info
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        file_size = os.path.getsize(pptx_path) / (1024 * 1024)  # MB
        
        return {
            "success": True,
            "presentation_path": pptx_path,
            "file_size_mb": round(file_size, 2),
            "instructions": {
                "step_1": "Go to https://gemini.google.com/",
                "step_2": "Click on 'Upload' or drag and drop the PowerPoint file",
                "step_3": "Ask Gemini to analyze the presentation, for example:",
                "example_prompts": [
                    "Analyze this PowerPoint presentation and provide feedback on visual design, content organization, and readability.",
                    "Review this presentation and suggest improvements for professional quality.",
                    "Evaluate the slides for clarity, visual balance, and audience appropriateness.",
                    "Extract key information from each slide and summarize the presentation."
                ],
                "step_4": "Gemini will analyze the slides visually and provide detailed feedback"
            },
            "alternative": "You can also use the analyze_presentation() method if you have Gemini API key set up"
        }


def analyze_presentation_vlm(pptx_path: str, 
                             api_key: Optional[str] = None,
                             backend: str = "auto",
                             model_name: Optional[str] = None,
                             generate_improved: bool = True) -> Dict[str, Any]:
    """
    Convenience function to analyze a PowerPoint presentation with VLM
    
    Args:
        pptx_path: Path to PowerPoint file
        api_key: Optional API key (for Gemini backend)
        backend: Backend to use ("auto", "local", "ollama", "gemini", "text")
        model_name: Model name for local/Ollama backends
        generate_improved: If True, generate improved slide content (default: True)
        
    Returns:
        Analysis results dictionary with improved slides
    """
    analyzer = VLMAnalyzer(api_key=api_key, backend=backend, model_name=model_name)
    return analyzer.analyze_presentation(pptx_path, generate_improved=generate_improved)

