"""
PowerPoint Presentation Builder
Creates PowerPoint files from ChatGPT-generated slide data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any, Optional
import os


def create_presentation_from_slides_data(slides_data: Dict[str, Any], 
                                        output_path: str) -> str:
    """
    Create PowerPoint presentation from ChatGPT-generated slides data
    
    Args:
        slides_data: Dictionary with title_slide and slides array
        output_path: Path where PPTX file should be saved
        
    Returns:
        Path to saved presentation file
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    if 'title_slide' in slides_data:
        title_slide_data = slides_data['title_slide']
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        title_shape = title_slide.shapes.title
        subtitle_shape = title_slide.placeholders[1]
        
        title_shape.text = title_slide_data.get('title', 'Presentation')
        subtitle_shape.text = title_slide_data.get('subtitle', '')
        
        # Format title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Format subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    else:
        # Default title slide if not provided
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = title_slide.shapes.title
        title_shape.text = "Presentation"
    
    # Create content slides
    slides = slides_data.get('slides', [])
    
    for slide_data in slides:
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        slide_title = slide_data.get('title', f"Slide {slide_data.get('slide_number', '?')}")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add content
        content = slide_data.get('content', [])
        content_y = Inches(1.3)
        content_height = Inches(5.5)
        content_width = Inches(9)
        
        content_box = slide.shapes.add_textbox(Inches(0.5), content_y, content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add content items
        if isinstance(content, list):
            for i, item in enumerate(content):
                if i > 0:
                    content_frame.add_paragraph()
                
                p = content_frame.paragraphs[i] if i < len(content_frame.paragraphs) else content_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(16)
                p.space_after = Pt(8)
                p.level = 0
                
                # Format bullet points
                if item.strip().startswith('•') or item.strip().startswith('-'):
                    p.level = 0
                elif item.strip().startswith('  ') or item.strip().startswith('\t'):
                    p.level = 1
        else:
            # Single text block
            p = content_frame.paragraphs[0]
            p.text = str(content)
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        
        # Format all paragraphs
        for paragraph in content_frame.paragraphs:
            paragraph.line_spacing = 1.2
            if paragraph.font.size is None:
                paragraph.font.size = Pt(16)
    
    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    
    # Save presentation
    prs.save(output_path)
    return output_path

